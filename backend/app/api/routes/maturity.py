"""
Legal AI OS — AI Maturity Assessment API

Governed API for law firm AI maturity assessment:
- Upload documents (bulk supported)
- Trigger maturity assessment
- View past assessments

Every route enforces client-scoping via get_current_user.
"""

from __future__ import annotations

import json
import traceback
from datetime import datetime, timezone
from uuid import UUID, uuid4

from fastapi import APIRouter, Depends, HTTPException, Query, UploadFile, File, Form

from app.auth import get_current_user, User
from app.database import get_supabase
from app.config import settings

router = APIRouter()


# ---------------------------------------------------------------------------
# Governance contract
# ---------------------------------------------------------------------------

@router.get("/health")
async def maturity_health():
    return {
        "function": "maturity-assessment",
        "status": "healthy",
        "version": "1.0.0",
        "capabilities": [
            "document_upload",
            "bulk_document_upload",
            "maturity_assessment",
            "dimension_scoring",
            "bottleneck_identification",
            "stage_gap_analysis",
        ],
    }


@router.get("/metrics")
async def maturity_metrics(user: User = Depends(get_current_user)):
    if not user.client_id:
        raise HTTPException(status_code=400, detail="No client association")

    supabase = get_supabase()
    assessments = (
        supabase.table("maturity_assessments")
        .select("id, document_count, cost_usd")
        .eq("client_id", str(user.client_id))
        .execute()
    )

    total_docs = sum(
        a.get("document_count", 0) for a in (assessments.data or [])
    )
    total_cost = sum(
        float(a.get("cost_usd", 0)) for a in (assessments.data or [])
    )

    return {
        "function": "maturity-assessment",
        "assessments_run": len(assessments.data or []),
        "total_documents_analyzed": total_docs,
        "total_cost_usd": round(total_cost, 6),
    }


@router.get("/targets")
async def maturity_targets():
    return {
        "function": "maturity-assessment",
        "targets": {
            "assessment_latency_seconds": "< 60",
            "dimensions_scored": 5,
            "levels": ["AI Aware", "AI Ready", "AI Capable", "AI Mature", "AI Native"],
            "max_documents_per_assessment": 100,
            "max_file_size_mb": 10,
        },
    }


# ---------------------------------------------------------------------------
# Assessments
# ---------------------------------------------------------------------------

@router.get("/assessments")
async def list_assessments(
    user: User = Depends(get_current_user),
    limit: int = Query(default=20, le=100),
):
    """List past maturity assessments for the user's client."""
    if not user.client_id:
        raise HTTPException(status_code=400, detail="No client association")

    result = (
        get_supabase()
        .table("maturity_assessments")
        .select(
            "id, version, overall_level, overall_level_label, "
            "bottleneck_dimension, document_count, cost_usd, created_at"
        )
        .eq("client_id", str(user.client_id))
        .order("created_at", desc=True)
        .limit(limit)
        .execute()
    )
    return result.data or []


@router.get("/assessments/{assessment_id}")
async def get_assessment(
    assessment_id: UUID,
    user: User = Depends(get_current_user),
):
    """Get a single maturity assessment with full detail."""
    if not user.client_id:
        raise HTTPException(status_code=400, detail="No client association")

    result = (
        get_supabase()
        .table("maturity_assessments")
        .select("*")
        .eq("id", str(assessment_id))
        .eq("client_id", str(user.client_id))
        .execute()
    )

    if not result.data:
        raise HTTPException(status_code=404, detail="Assessment not found")

    data = result.data[0]
    # Parse JSONB fields
    if isinstance(data.get("dimensions"), str):
        try:
            data["dimensions"] = json.loads(data["dimensions"])
        except (json.JSONDecodeError, TypeError):
            data["dimensions"] = []
    if isinstance(data.get("stage_gaps"), str):
        try:
            data["stage_gaps"] = json.loads(data["stage_gaps"])
        except (json.JSONDecodeError, TypeError):
            data["stage_gaps"] = []

    return data


@router.post("/assessments", status_code=201)
async def run_assessment(
    data: dict,
    user: User = Depends(get_current_user),
):
    """
    Trigger a new maturity assessment.

    Body: {"document_ids": ["uuid1", "uuid2", ...]}

    The assessment runs synchronously (typically < 30s for up to 50 docs).
    For very large document sets (> 100K chars), a two-pass summarization
    is applied first.
    """
    if not user.client_id:
        raise HTTPException(status_code=400, detail="No client association")

    document_ids = data.get("document_ids", [])
    if not document_ids:
        raise HTTPException(
            status_code=400,
            detail="At least one document_id is required",
        )

    if len(document_ids) > 100:
        raise HTTPException(
            status_code=400,
            detail="Maximum 100 documents per assessment",
        )

    try:
        from app.services.maturity import MaturityService
        service = MaturityService()
        assessment = await service.run_assessment(
            client_id=user.client_id,
            user_id=user.id,
            document_ids=[UUID(did) for did in document_ids],
        )

        # Parse JSONB for response
        if isinstance(assessment.get("dimensions"), str):
            try:
                assessment["dimensions"] = json.loads(assessment["dimensions"])
            except (json.JSONDecodeError, TypeError):
                assessment["dimensions"] = []
        if isinstance(assessment.get("stage_gaps"), str):
            try:
                assessment["stage_gaps"] = json.loads(assessment["stage_gaps"])
            except (json.JSONDecodeError, TypeError):
                assessment["stage_gaps"] = []

        return assessment

    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))
    except Exception as e:
        traceback.print_exc()
        raise HTTPException(
            status_code=500, detail=f"Assessment failed: {str(e)}"
        )


# ---------------------------------------------------------------------------
# Document management
# ---------------------------------------------------------------------------

SUPPORTED_EXTENSIONS = {
    ".pdf", ".docx", ".doc", ".txt", ".md",
    ".csv", ".xlsx", ".pptx", ".rtf", ".html", ".xml",
}
MAX_FILE_SIZE = 10 * 1024 * 1024  # 10 MB


@router.get("/documents")
async def list_documents(
    user: User = Depends(get_current_user),
    limit: int = Query(default=100, le=500),
    offset: int = 0,
):
    """List documents uploaded for maturity assessment."""
    if not user.client_id:
        raise HTTPException(status_code=400, detail="No client association")

    result = (
        get_supabase()
        .table("knowledge_documents")
        .select(
            "id, title, document_type, source_file, "
            "chunk_count, is_active, created_at"
        )
        .eq("client_id", str(user.client_id))
        .order("created_at", desc=True)
        .limit(limit)
        .offset(offset)
        .execute()
    )
    return result.data or []


@router.post("/documents/upload", status_code=201)
async def upload_document(
    file: UploadFile = File(...),
    user: User = Depends(get_current_user),
):
    """
    Upload a single document for maturity assessment.

    Supported: PDF, DOCX, DOC, TXT, MD, CSV, XLSX, PPTX, RTF, HTML, XML.
    Max 10 MB per file.
    """
    return await _handle_upload(file, user)


@router.post("/documents/upload-bulk", status_code=201)
async def upload_documents_bulk(
    files: list[UploadFile] = File(...),
    user: User = Depends(get_current_user),
):
    """
    Upload multiple documents at once. Max 20 files per call.
    """
    if len(files) > 20:
        raise HTTPException(
            status_code=400, detail="Maximum 20 files per bulk upload"
        )

    results = []
    errors = []
    for file in files:
        try:
            result = await _handle_upload(file, user)
            results.append(result)
        except HTTPException as e:
            errors.append({"filename": file.filename, "error": str(e.detail)})
        except Exception as e:
            errors.append({"filename": file.filename, "error": str(e)})

    return {
        "uploaded": results,
        "errors": errors,
        "total": len(files),
        "succeeded": len(results),
        "failed": len(errors),
    }


@router.delete("/documents/{document_id}")
async def delete_document(
    document_id: UUID,
    user: User = Depends(get_current_user),
):
    """Soft-delete a document (set is_active=false)."""
    if not user.client_id:
        raise HTTPException(status_code=400, detail="No client association")

    result = (
        get_supabase()
        .table("knowledge_documents")
        .update({"is_active": False})
        .eq("id", str(document_id))
        .eq("client_id", str(user.client_id))
        .execute()
    )
    return {"status": "deactivated"}


# ---------------------------------------------------------------------------
# Upload helpers
# ---------------------------------------------------------------------------

def _extract_text(file_bytes: bytes, filename: str) -> str:
    """Extract text from uploaded file bytes."""
    import io
    import os

    ext = os.path.splitext(filename)[1].lower()

    if ext in (".txt", ".md", ".html", ".xml", ".csv", ".rtf"):
        try:
            return file_bytes.decode("utf-8")
        except UnicodeDecodeError:
            try:
                return file_bytes.decode("latin-1")
            except Exception:
                return f"[Binary or undecodable file: {filename}]"

    if ext == ".pdf":
        try:
            import pypdf
            reader = pypdf.PdfReader(io.BytesIO(file_bytes))
            pages = []
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    pages.append(text)
            return "\n\n".join(pages) if pages else f"[No extractable text in PDF: {filename}]"
        except ImportError:
            return f"[PDF extraction unavailable: {filename}]"
        except Exception as e:
            return f"[PDF extraction error: {filename} — {e}]"

    if ext in (".docx", ".doc"):
        try:
            import docx
            doc = docx.Document(io.BytesIO(file_bytes))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n".join(paragraphs) if paragraphs else f"[Empty DOCX: {filename}]"
        except ImportError:
            return f"[DOCX extraction unavailable: {filename}]"
        except Exception as e:
            return f"[DOCX extraction error: {filename} — {e}]"

    if ext == ".xlsx":
        try:
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True)
            rows = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows.append(f"Sheet: {sheet_name}")
                for row in ws.iter_rows(values_only=True):
                    rows.append("\t".join(str(c) if c is not None else "" for c in row))
            return "\n".join(rows) if rows else f"[Empty XLSX: {filename}]"
        except ImportError:
            return f"[XLSX extraction unavailable: {filename}]"
        except Exception as e:
            return f"[XLSX extraction error: {filename} — {e}]"

    if ext == ".pptx":
        try:
            import pptx
            prs = pptx.Presentation(io.BytesIO(file_bytes))
            slides = []
            for slide in prs.slides:
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        texts.append(shape.text_frame.text)
                slides.append("\n".join(texts))
            return "\n\n---\n\n".join(slides) if slides else f"[Empty PPTX: {filename}]"
        except ImportError:
            return f"[PPTX extraction unavailable: {filename}]"
        except Exception as e:
            return f"[PPTX extraction error: {filename} — {e}]"

    return f"[Unsupported file type: {filename}]"


async def _handle_upload(file: UploadFile, user: User) -> dict:
    """Validate, extract, and store a single uploaded document."""
    import os

    if not user.client_id:
        raise HTTPException(status_code=400, detail="No client association")

    # Validate extension
    ext = os.path.splitext(file.filename or "")[1].lower()
    if ext not in SUPPORTED_EXTENSIONS:
        raise HTTPException(
            status_code=400,
            detail=f"Unsupported file type: {ext}. Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}",
        )

    # Read and validate size
    contents = await file.read()
    if len(contents) > MAX_FILE_SIZE:
        raise HTTPException(
            status_code=400,
            detail=f"File too large: {len(contents)} bytes. Max: {MAX_FILE_SIZE} bytes.",
        )

    # Extract text
    text = _extract_text(contents, file.filename or "unknown")

    # Determine document type
    doc_type = _doc_type_from_ext(ext)

    # Generate embedding if Voyage is available
    embedding = await _embed_text(text[:8000])

    # Store in knowledge_documents
    title = os.path.splitext(file.filename or "Untitled")[0]
    chunk_count = max(1, len(text) // 800)

    result = (
        get_supabase()
        .table("knowledge_documents")
        .insert({
            "client_id": str(user.client_id),
            "title": title,
            "document_type": doc_type,
            "source_file": file.filename,
            "raw_text": text,
            "chunk_count": chunk_count,
            "metadata": {
                "source": "maturity_upload",
                "file_size_bytes": len(contents),
                "extension": ext,
            },
            "embedding": embedding,
            "created_by": str(user.id),
        })
        .execute()
    )

    if not result.data:
        raise HTTPException(status_code=500, detail="Failed to store document")

    doc = result.data[0]
    return {
        "id": doc["id"],
        "title": doc["title"],
        "document_type": doc["document_type"],
        "source_file": doc.get("source_file"),
        "chunk_count": doc.get("chunk_count", chunk_count),
        "created_at": doc.get("created_at"),
    }


def _doc_type_from_ext(ext: str) -> str:
    """Map file extension to document type."""
    mapping = {
        ".pdf": "document",
        ".docx": "document",
        ".doc": "document",
        ".txt": "memo",
        ".md": "memo",
        ".csv": "data",
        ".xlsx": "data",
        ".pptx": "presentation",
        ".rtf": "document",
        ".html": "web",
        ".xml": "data",
    }
    return mapping.get(ext, "document")


async def _embed_text(text: str) -> list[float] | None:
    """Generate embedding via Voyage AI, if configured."""
    import asyncio
    try:
        return await asyncio.to_thread(_embed_sync, text)
    except Exception:
        return None


def _embed_sync(text: str) -> list[float] | None:
    """Synchronous Voyage AI embedding call."""
    if not settings.voyage_api_key:
        return None
    try:
        import voyageai
        vo = voyageai.Client()
        result = vo.embed(
            text, model=settings.voyage_model, input_type="document"
        )
        return result.embeddings[0] if result.embeddings else None
    except Exception:
        return None
