from logger_config import get_logger
logger = get_logger(__name__)

"""
Document processing module for SuperAssistant
Handles text extraction, chunking, and embedding generation
"""
import os
from typing import List, Dict
from io import BytesIO
from datetime import datetime, timezone
import voyageai
from docx import Document
from database import get_supabase
from dotenv import load_dotenv
from config.constants import TEXT_CHUNKING, EMBEDDING, SEARCH, DATABASE
from cache import get_cached_search_results, cache_search_results
from api.utils.error_handler import (
    DocumentProcessingError,
    TextExtractionError,
    UnsupportedFileTypeError,
    EmbeddingGenerationError,
    wrap_external_service_error
)
from api.utils.retry import retry_with_backoff, retry_on_rate_limit
try:
    from openpyxl import load_workbook
    XLSX_SUPPORT = True
except ImportError:
    XLSX_SUPPORT = False
    logger.warning("openpyxl not installed - XLSX support disabled")

try:
    from pptx import Presentation
    PPTX_SUPPORT = True
except ImportError:
    PPTX_SUPPORT = False
    logger.warning("python-pptx not installed - PPTX support disabled")

load_dotenv()

# Initialize clients
VOYAGE_API_KEY = os.environ.get("VOYAGE_API_KEY")
SUPABASE_URL = os.environ.get("SUPABASE_URL")

vo = voyageai.Client(api_key=VOYAGE_API_KEY)
supabase = get_supabase()


def extract_text_from_file(file_data: bytes, filename: str) -> str:
    """
    Extract text from different file types

    Args:
        file_data: Raw file bytes
        filename: Name of the file (used to determine type)

    Returns:
        Extracted text content
    """
    file_ext = filename.lower().split('.')[-1]

    if file_ext == 'docx':
        # Extract text from Word document
        try:
            doc = Document(BytesIO(file_data))
            text = '\n\n'.join([paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()])
            if not text.strip():
                raise ValueError("Document appears to be empty or contains no readable text")
            return text
        except Exception as e:
            logger.error(f"❌ Error extracting text from .docx file: {str(e)}")
            raise ValueError(f"Failed to extract text from Word document: {str(e)}")

    elif file_ext == 'xlsx':
        # Extract text from Excel spreadsheet
        if not XLSX_SUPPORT:
            raise ValueError("XLSX support not available - install openpyxl: pip install openpyxl")
        try:
            workbook = load_workbook(BytesIO(file_data), data_only=True)
            text_parts = []
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_parts.append(f"=== Sheet: {sheet_name} ===\n")
                for row in sheet.iter_rows(values_only=True):
                    row_text = '\t'.join([str(cell) if cell is not None else '' for cell in row])
                    if row_text.strip():
                        text_parts.append(row_text)
            text = '\n'.join(text_parts)
            if not text.strip():
                raise ValueError("Spreadsheet appears to be empty")
            return text
        except Exception as e:
            logger.error(f"❌ Error extracting text from .xlsx file: {str(e)}")
            raise ValueError(f"Failed to extract text from Excel file: {str(e)}")

    elif file_ext == 'pptx':
        # Extract text from PowerPoint presentation
        if not PPTX_SUPPORT:
            raise ValueError("PPTX support not available - install python-pptx: pip install python-pptx")
        try:
            prs = Presentation(BytesIO(file_data))
            text_parts = []
            for slide_num, slide in enumerate(prs.slides, 1):
                text_parts.append(f"=== Slide {slide_num} ===\n")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)
            text = '\n\n'.join(text_parts)
            if not text.strip():
                raise ValueError("Presentation appears to be empty")
            return text
        except Exception as e:
            logger.error(f"❌ Error extracting text from .pptx file: {str(e)}")
            raise ValueError(f"Failed to extract text from PowerPoint file: {str(e)}")

    elif file_ext == 'pdf':
        # Extract text from PDF
        try:
            import PyPDF2
            from io import BytesIO
            pdf_reader = PyPDF2.PdfReader(BytesIO(file_data))
            text_parts = []
            for page_num, page in enumerate(pdf_reader.pages, 1):
                text = page.extract_text()
                if text.strip():
                    text_parts.append(f"=== Page {page_num} ===\n{text}")
            full_text = '\n\n'.join(text_parts)
            if not full_text.strip():
                raise ValueError("PDF appears to be empty or contains no extractable text")
            return full_text
        except ImportError:
            raise ValueError("PDF support not available - install PyPDF2: pip install PyPDF2")
        except Exception as e:
            logger.error(f"❌ Error extracting text from .pdf file: {str(e)}")
            raise ValueError(f"Failed to extract text from PDF: {str(e)}")

    elif file_ext == 'txt':
        # Decode text file
        try:
            return file_data.decode('utf-8')
        except UnicodeDecodeError as e:
            raise ValueError(f"Failed to decode text file: {str(e)}")

    else:
        # Try to decode as text (fallback)
        try:
            return file_data.decode('utf-8')
        except UnicodeDecodeError:
            raise ValueError(f"Unsupported file type: .{file_ext}. Supported formats: .pdf, .txt, .docx, .xlsx, .pptx")


def chunk_text(
    text: str,
    chunk_size: int = TEXT_CHUNKING.DEFAULT_CHUNK_SIZE,
    overlap: int = TEXT_CHUNKING.DEFAULT_OVERLAP
) -> List[Dict]:
    """
    Split text into overlapping chunks

    Args:
        text: The text to chunk
        chunk_size: Target size of each chunk in characters
        overlap: Number of characters to overlap between chunks

    Returns:
        List of dicts with 'content' and 'chunk_index'
    """
    chunks = []
    start = 0
    chunk_index = 0

    while start < len(text):
        # Get chunk
        end = start + chunk_size
        chunk_text = text[start:end]

        # If not the last chunk, try to break at a sentence or word boundary
        if end < len(text):
            # Look for sentence ending
            last_period = chunk_text.rfind('.')
            last_question = chunk_text.rfind('?')
            last_exclamation = chunk_text.rfind('!')

            boundary = max(last_period, last_question, last_exclamation)

            # If found a sentence boundary in the last X% of chunk, use it
            if boundary > chunk_size * TEXT_CHUNKING.SENTENCE_BOUNDARY_THRESHOLD:
                end = start + boundary + 1
                chunk_text = text[start:end]

        # Only add non-empty chunks
        if chunk_text.strip():
            chunks.append({
                'content': chunk_text.strip(),
                'chunk_index': chunk_index
            })
            chunk_index += 1

        # Move start position (with overlap)
        start = end - overlap

    return chunks


@retry_with_backoff(
    max_retries=3,
    initial_delay=2.0,
    max_delay=30.0,
    exponential_base=2.0
)
def _call_voyage_embed_api(batch: List[str], model: str, input_type: str) -> List[List[float]]:
    """
    Internal function to call Voyage AI API with retry logic.

    This function is decorated with retry logic to handle transient failures
    such as network issues, rate limits, and temporary API unavailability.

    Args:
        batch: List of texts to embed
        model: Model name to use
        input_type: Type of input (document or query)

    Returns:
        List of embedding vectors

    Raises:
        EmbeddingGenerationError: If all retry attempts fail
    """
    try:
        result = vo.embed(batch, model=model, input_type=input_type)
        return result.embeddings
    except Exception as e:
        # Wrap external service errors with our custom exception
        raise wrap_external_service_error(e, "Voyage AI", "embedding generation")


def generate_embeddings(
    texts: List[str],
    batch_size: int = EMBEDDING.DEFAULT_BATCH_SIZE,
    input_type: str = EMBEDDING.INPUT_TYPE_DOCUMENT
) -> List[List[float]]:
    """
    Generate embeddings for a list of texts using Voyage AI
    Handles batching to respect VoyageAI API limits:
    - Max 1000 texts per batch
    - Max 320,000 tokens per batch (~1.28M characters assuming 4 chars/token)

    Args:
        texts: List of text strings to embed
        batch_size: Number of texts to process per batch (default from EMBEDDING config)
        input_type: Type of input - "document" for indexing, "query" for search queries

    Returns:
        List of embedding vectors (1024 dimensions each)
    """
    all_embeddings = []
    total_batches = (len(texts) + batch_size - 1) // batch_size

    logger.info(f"   Processing {len(texts)} texts in {total_batches} batches (batch_size={batch_size}, input_type={input_type})...")

    try:
        for i in range(0, len(texts), batch_size):
            batch_num = (i // batch_size) + 1
            batch = texts[i:i + batch_size]

            logger.info(f"      Processing batch {batch_num}/{total_batches} ({len(batch)} texts)...")

            # Call Voyage AI API with retry logic
            embeddings = _call_voyage_embed_api(
                batch=batch,
                model=EMBEDDING.MODEL_NAME,
                input_type=input_type
            )
            all_embeddings.extend(embeddings)

        logger.info(f"   ✓ Successfully generated {len(all_embeddings)} embeddings")
        return all_embeddings

    except Exception as e:
        logger.error(f"❌ Error generating embeddings: {str(e)}")
        raise EmbeddingGenerationError(
            f"Failed to generate embeddings after retries: {str(e)}",
            details={"batch_count": total_batches, "text_count": len(texts)}
        )


def process_document(document_id: str) -> Dict:
    """
    Process a document: extract text, chunk, embed, store

    Args:
        document_id: UUID of the document to process

    Returns:
        Processing result summary
    """
    logger.info(f"\n📄 Processing document: {document_id}")

    try:
        # Get document metadata
        doc_result = supabase.table('documents').select('*').eq('id', document_id).execute()

        if not doc_result.data:
            raise ValueError(f"Document {document_id} not found")

        document = doc_result.data[0]
        logger.info(f"   Filename: {document['filename']}")
        logger.info(f"   Client ID: {document['client_id']}")

        # Set status to processing
        supabase.table('documents').update({
            'processing_status': 'processing'
        }).eq('id', document_id).execute()

        # Download file from storage
        storage_path = document['storage_url'].split('/documents/')[-1]
        file_data = supabase.storage.from_('documents').download(storage_path)

        # Extract text based on file type
        text = extract_text_from_file(file_data, document['filename'])
        logger.info(f"   Text length: {len(text)} characters")

        # Chunk text
        chunks = chunk_text(text, chunk_size=800, overlap=200)
        logger.info(f"   Created {len(chunks)} chunks")

        # Extract text content for embedding
        chunk_texts = [chunk['content'] for chunk in chunks]

        # Generate embeddings
        logger.info(f"   Generating embeddings...")
        embeddings = generate_embeddings(chunk_texts)
        logger.info(f"   ✓ Generated {len(embeddings)} embeddings")

        # Store chunks with embeddings in database
        logger.info(f"   Storing chunks in database...")

        # FIX: Batch insert instead of individual inserts (10-12x performance improvement)
        # Collect all chunk data first
        chunks_to_insert = []
        for i, (chunk, embedding) in enumerate(zip(chunks, embeddings)):
            # Sanitize content: remove null bytes which PostgreSQL cannot handle
            sanitized_content = chunk['content'].replace('\x00', '')

            chunk_data = {
                'document_id': document_id,
                'client_id': document['client_id'],
                'content': sanitized_content,
                'embedding': embedding,
                'chunk_index': chunk['chunk_index'],
                'metadata': {
                    'filename': document['filename'],
                    'chunk_size': len(sanitized_content)
                }
            }
            chunks_to_insert.append(chunk_data)

        # Batch insert with PostgreSQL limit handling
        BATCH_SIZE = DATABASE.BATCH_INSERT_SIZE
        stored_count = 0

        for batch_start in range(0, len(chunks_to_insert), BATCH_SIZE):
            batch = chunks_to_insert[batch_start:batch_start + BATCH_SIZE]
            supabase.table('document_chunks').insert(batch).execute()
            stored_count += len(batch)

        logger.info(f"   ✓ Stored {stored_count} chunks in {(len(chunks_to_insert) + BATCH_SIZE - 1) // BATCH_SIZE} batch(es)")

        # Mark document as processed
        supabase.table('documents').update({
            'processed': True,
            'processed_at': datetime.now(timezone.utc).isoformat(),
            'processing_status': 'completed',
            'processing_error': None
        }).eq('id', document_id).execute()

        logger.info(f"   ✓ Marked document as processed")

        return {
            'document_id': document_id,
            'filename': document['filename'],
            'chunks_created': len(chunks),
            'chunks_stored': stored_count,
            'text_length': len(text),
            'status': 'completed'
        }

    except DocumentProcessingError as e:
        # Handle known document processing errors
        error_message = e.message if hasattr(e, 'message') else str(e)
        logger.error(
            f"   ❌ Document processing error: {error_message}",
            extra={
                "document_id": document_id,
                "error_type": e.__class__.__name__,
                "error_code": e.error_code if hasattr(e, 'error_code') else None
            }
        )

        # Update document with error status
        try:
            supabase.table('documents').update({
                'processing_status': 'failed',
                'processing_error': error_message,
                'processed': True  # Mark as processed even on failure so UI knows processing is done
            }).eq('id', document_id).execute()
            logger.info(f"   ✓ Marked document as failed")
        except Exception as update_error:
            logger.error(f"   ⚠️  Failed to update error status: {update_error}")

        # Don't re-raise - background task should not fail
        return {
            'document_id': document_id,
            'status': 'failed',
            'error': error_message
        }

    except Exception as e:
        # Handle unexpected errors
        error_message = f"Unexpected error: {str(e)}"
        logger.exception(
            f"   ❌ Unexpected processing error",
            extra={"document_id": document_id},
            exc_info=True
        )

        # Update document with error status
        try:
            supabase.table('documents').update({
                'processing_status': 'failed',
                'processing_error': error_message,
                'processed': True  # Mark as processed even on failure so UI knows processing is done
            }).eq('id', document_id).execute()
            logger.info(f"   ✓ Marked document as failed")
        except Exception as update_error:
            logger.error(f"   ⚠️  Failed to update error status: {update_error}")

        # Don't re-raise - background task should not fail
        return {
            'document_id': document_id,
            'status': 'failed',
            'error': error_message
        }


def process_conversation_to_kb(conversation_id: str, user_id: str) -> Dict:
    """
    Process a conversation into chunks and add to knowledge base

    Args:
        conversation_id: UUID of the conversation
        user_id: UUID of the user (for access control)

    Returns:
        Processing result summary
    """
    logger.info(f"\n💬 Processing conversation to KB: {conversation_id}")

    # Get conversation metadata and verify ownership
    conv_result = supabase.table('conversations')\
        .select('*')\
        .eq('id', conversation_id)\
        .eq('user_id', user_id)\
        .execute()

    if not conv_result.data:
        raise ValueError(f"Conversation {conversation_id} not found or access denied")

    conversation = conv_result.data[0]
    logger.info(f"   Title: {conversation.get('title', 'Untitled')}")
    logger.info(f"   User ID: {conversation['user_id']}")

    # Check if conversation is already in knowledge base
    if conversation.get('in_knowledge_base'):
        logger.warning(f"   ⚠️  Conversation already in knowledge base")
        raise ValueError(f"Conversation '{conversation.get('title', 'Untitled')}' is already in your knowledge base")

    # Get all messages from the conversation
    messages_result = supabase.table('messages')\
        .select('role, content, timestamp')\
        .eq('conversation_id', conversation_id)\
        .order('timestamp')\
        .execute()

    if not messages_result.data:
        raise ValueError(f"No messages found in conversation {conversation_id}")

    messages = messages_result.data
    logger.info(f"   Found {len(messages)} messages")

    # Combine messages into conversation text with formatting
    conversation_text = f"Conversation: {conversation.get('title', 'Untitled')}\n\n"
    for msg in messages:
        role_label = "User" if msg['role'] == 'user' else "Assistant"
        conversation_text += f"{role_label}: {msg['content']}\n\n"

    logger.info(f"   Total conversation length: {len(conversation_text)} characters")

    # Chunk the conversation
    chunks = chunk_text(conversation_text, chunk_size=800, overlap=200)
    logger.info(f"   Created {len(chunks)} chunks")

    # Extract text content for embedding
    chunk_texts = [chunk['content'] for chunk in chunks]

    # Generate embeddings
    logger.info(f"   Generating embeddings...")
    embeddings = generate_embeddings(chunk_texts)
    logger.info(f"   ✓ Generated {len(embeddings)} embeddings")

    # Calculate embedding cost (Voyage AI voyage-3: ~$0.02 per 1M tokens)
    # Rough estimate: 4 characters ≈ 1 token for English
    total_chars = sum(len(text) for text in chunk_texts)
    estimated_tokens = total_chars / 4
    embedding_cost_usd = (estimated_tokens / 1_000_000) * 0.02
    logger.info(f"   💰 Embedding cost: ~${embedding_cost_usd:.6f} ({estimated_tokens:.0f} tokens)")

    # Get client_id from user
    user_result = supabase.table('users')\
        .select('client_id')\
        .eq('id', user_id)\
        .single()\
        .execute()

    client_id = user_result.data['client_id']

    # Store chunks with embeddings in database with transaction safety
    logger.info(f"   Storing chunks in database...")

    try:
        # FIX: Batch insert instead of individual inserts (10-12x performance improvement)
        # Collect all chunk data first
        chunks_to_insert = []
        for i, (chunk, embedding) in enumerate(zip(chunks, embeddings)):
            # Sanitize content: remove null bytes which PostgreSQL cannot handle
            sanitized_content = chunk['content'].replace('\x00', '')

            chunk_data = {
                'conversation_id': conversation_id,
                'client_id': client_id,
                'source_type': 'conversation',
                'content': sanitized_content,
                'embedding': embedding,
                'chunk_index': chunk['chunk_index'],
                'metadata': {
                    'conversation_title': conversation.get('title', 'Untitled'),
                    'message_count': len(messages),
                    'chunk_size': len(sanitized_content)
                }
            }
            chunks_to_insert.append(chunk_data)

        # Batch insert with PostgreSQL limit handling
        BATCH_SIZE = DATABASE.BATCH_INSERT_SIZE
        stored_count = 0

        for batch_start in range(0, len(chunks_to_insert), BATCH_SIZE):
            batch = chunks_to_insert[batch_start:batch_start + BATCH_SIZE]
            supabase.table('document_chunks').insert(batch).execute()
            stored_count += len(batch)

        logger.info(f"   ✓ Stored {stored_count} chunks in {(len(chunks_to_insert) + BATCH_SIZE - 1) // BATCH_SIZE} batch(es)")

        # Only mark conversation as in knowledge base if ALL chunks succeeded
        supabase.table('conversations').update({
            'in_knowledge_base': True,
            'added_to_kb_at': 'now()'
        }).eq('id', conversation_id).execute()

        logger.info(f"   ✓ Marked conversation as in knowledge base")

        return {
            'conversation_id': conversation_id,
            'title': conversation.get('title', 'Untitled'),
            'chunks_created': len(chunks),
            'chunks_stored': stored_count,
            'message_count': len(messages),
            'embedding_cost_usd': round(embedding_cost_usd, 6),
            'estimated_tokens': int(estimated_tokens),
            'status': 'completed'
        }

    except Exception as e:
        # Rollback: Clean up any partially inserted chunks
        logger.error(f"   ❌ Error during KB processing: {str(e)}")
        logger.info(f"   🔄 Rolling back {stored_count} chunks...")

        try:
            supabase.table('document_chunks')\
                .delete()\
                .eq('conversation_id', conversation_id)\
                .eq('source_type', 'conversation')\
                .execute()
            logger.info(f"   ✓ Rollback complete")
        except Exception as rollback_error:
            logger.error(f"   ⚠️  Rollback failed: {rollback_error}")

        # Re-raise the original error with context
        raise ValueError(f"Failed to add conversation to knowledge base: {str(e)}")


def remove_conversation_from_kb(conversation_id: str, user_id: str) -> Dict:
    """
    Remove a conversation from the knowledge base

    Args:
        conversation_id: UUID of the conversation
        user_id: UUID of the user (for access control)

    Returns:
        Removal result summary
    """
    logger.info(f"\n🗑️  Removing conversation from KB: {conversation_id}")

    # Verify ownership
    conv_result = supabase.table('conversations')\
        .select('id')\
        .eq('id', conversation_id)\
        .eq('user_id', user_id)\
        .execute()

    if not conv_result.data:
        raise ValueError(f"Conversation {conversation_id} not found or access denied")

    # Delete all chunks for this conversation
    delete_result = supabase.table('document_chunks')\
        .delete()\
        .eq('conversation_id', conversation_id)\
        .execute()

    chunks_deleted = len(delete_result.data) if delete_result.data else 0
    logger.info(f"   ✓ Deleted {chunks_deleted} chunks")

    # Mark conversation as not in knowledge base
    supabase.table('conversations').update({
        'in_knowledge_base': False,
        'added_to_kb_at': None
    }).eq('id', conversation_id).execute()

    logger.info(f"   ✓ Marked conversation as removed from knowledge base")

    return {
        'conversation_id': conversation_id,
        'chunks_deleted': chunks_deleted,
        'status': 'removed'
    }


def preprocess_query(query: str) -> str:
    """
    Preprocess search query to improve relevance by removing stop words
    and focusing on key terms

    Args:
        query: The raw search query

    Returns:
        Preprocessed query with stop words removed
    """
    # Common stop words that dilute semantic meaning
    stop_words = {
        'what', 'are', 'is', 'the', 'a', 'an', 'and', 'or', 'but', 'in', 'on', 'at',
        'to', 'for', 'of', 'with', 'by', 'from', 'as', 'this', 'that', 'these', 'those',
        'i', 'you', 'we', 'they', 'it', 'my', 'your', 'our', 'their', 'can', 'could',
        'would', 'should', 'will', 'do', 'does', 'did', 'have', 'has', 'had', 'be', 'been',
        'am', 'was', 'were', 'me', 'us', 'them', 'there', 'here', 'where', 'when', 'how',
        'why', 'which', 'who', 'whom', 'whose', 'tell', 'show', 'give', 'find', 'get',
        # Meta-terms that don't add semantic value to queries
        'any', 'some', 'about', 'regarding', 'information', 'info', 'data', 'details',
        'know', 'help', 'need', 'want', 'please'
    }

    # Split query into words and filter out stop words
    words = query.lower().split()
    key_terms = [word.strip('.,?!:;') for word in words if word.lower() not in stop_words]

    # If all words were stop words, return original query to avoid empty search
    if not key_terms:
        return query

    # Reconstruct query from key terms
    preprocessed = ' '.join(key_terms)

    logger.info(f"   Query preprocessing:")
    logger.info(f"      Original: '{query}'")
    logger.info(f"      Preprocessed: '{preprocessed}'")

    return preprocessed


def detect_query_type(query: str) -> str:
    """
    Detect the type of query to determine appropriate search strategy

    Args:
        query: The search query

    Returns:
        Query type: 'factual' for direct questions, 'exploratory' for general queries
    """
    query_lower = query.lower().strip()

    # Direct question words indicate a need for specific factual answers
    question_words = ['what', 'who', 'where', 'when', 'how', 'why', 'which', 'whose']

    # Check if query starts with a question word
    for qword in question_words:
        if query_lower.startswith(qword + ' '):
            return 'factual'

    # Questions with "is/are/do/does" patterns
    question_patterns = [
        'is there', 'are there', 'is the', 'are the', 'do we', 'does the',
        'can you', 'could you', 'will the', 'would the',
        'do you', 'did you', 'have you', 'has there', 'does it'
    ]

    for pattern in question_patterns:
        if query_lower.startswith(pattern):
            return 'factual'

    # Default to exploratory for general statements or broad queries
    return 'exploratory'


def search_similar_chunks(query: str, client_id: str, limit: int = 5, include_conversations: bool = True, min_similarity: float = 0.0) -> List[Dict]:
    """
    Search for document chunks similar to a query

    Uses Redis cache for identical queries (1-hour TTL) to avoid expensive
    vector similarity searches.

    Args:
        query: The search query
        client_id: Client ID to filter results
        limit: Maximum number of results
        include_conversations: Whether to include conversation chunks (default: True)
        min_similarity: Minimum similarity score to include results (default: 0.0)

    Returns:
        List of similar chunks with similarity scores
    """
    logger.info(f"\n🔍 Searching for: '{query}'")
    logger.info(f"   Client: {client_id}")
    logger.info(f"   Include conversations: {include_conversations}")
    logger.info(f"   Min similarity: {min_similarity}")

    # Try to get from cache first (1-hour TTL)
    # Cache key includes all parameters to ensure exact match
    cache_key_suffix = f"{query}:{limit}:{include_conversations}:{min_similarity}"
    cached_results = get_cached_search_results(cache_key_suffix, client_id)
    if cached_results is not None:
        logger.info(f"   ✅ Search results loaded from cache ({len(cached_results)} results)")
        return cached_results

    logger.debug("   📋 Cache miss - performing vector search")

    # Detect query type for adaptive strategy
    query_type = detect_query_type(query)
    logger.info(f"   Query type: {query_type}")

    # Preprocess query to extract key terms
    preprocessed_query = preprocess_query(query)

    # Adjust minimum similarity based on query type
    # Factual questions need higher precision, exploratory queries can be more permissive
    if min_similarity == 0.0:  # Only adjust if not explicitly set
        if query_type == 'factual':
            min_similarity = SEARCH.SIMILARITY_THRESHOLDS['factual']
            logger.info(f"   Adjusted min_similarity to {min_similarity} for factual query")
        else:
            min_similarity = SEARCH.SIMILARITY_THRESHOLDS['exploratory']

    # Generate embedding for query using "query" input type for better search relevance
    query_embedding = generate_embeddings([preprocessed_query], input_type=EMBEDDING.INPUT_TYPE_QUERY)[0]

    # Call vector search function
    result = supabase.rpc(
        'match_document_chunks',
        {
            'query_embedding': query_embedding,
            'match_count': limit,
            'filter_client_id': client_id
        }
    ).execute()

    chunks = result.data
    logger.info(f"   Found {len(chunks)} results")

    # Filter out conversation chunks if requested
    if not include_conversations:
        chunks = [c for c in chunks if c.get('source_type') != 'conversation']
        logger.info(f"   After filtering: {len(chunks)} document-only results")

    # Filter by minimum similarity score
    if min_similarity > 0.0:
        chunks = [c for c in chunks if c.get('similarity', 0.0) >= min_similarity]
        logger.info(f"   After similarity filtering (>={min_similarity}): {len(chunks)} results")

    # Cache results for 1 hour (3600 seconds)
    cache_search_results(cache_key_suffix, client_id, chunks, ttl=3600)

    return chunks
